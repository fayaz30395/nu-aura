from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("docs/presentations/ai-for-modern-teams-real-world-usage-premium.pptx")
LOGO_PATH = Path("docs/assets/NU logo.png")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLORS = {
    "bg": RGBColor(7, 12, 24),
    "bg_soft": RGBColor(14, 24, 42),
    "panel": RGBColor(19, 32, 54),
    "panel_2": RGBColor(25, 42, 71),
    "blue": RGBColor(41, 121, 255),
    "cyan": RGBColor(34, 211, 238),
    "teal": RGBColor(20, 184, 166),
    "green": RGBColor(74, 222, 128),
    "amber": RGBColor(251, 191, 36),
    "rose": RGBColor(251, 113, 133),
    "white": RGBColor(245, 247, 250),
    "text": RGBColor(226, 232, 240),
    "muted": RGBColor(148, 163, 184),
    "line": RGBColor(51, 65, 85),
}


SLIDES = [
    {
        "layout": "title",
        "title": "AI for Modern Teams",
        "subtitle": "Real-world usage. Daily workflows. Immediate adoption.",
        "kicker": "LIVE WORKSHOP",
        "note": "Set the tone as practical and current. This is not an AI overview deck; it is a field guide for actual team usage.",
    },
    {
        "layout": "statement",
        "title": "The Reality Shift",
        "statement": "Teams using AI well are already shipping faster.",
        "bullets": [
            "AI is now embedded in coding, planning, writing, and analysis",
            "The gap is not hype versus no hype; it is workflow maturity versus lag",
            "High performers use AI to compress cycle time without lowering standards",
        ],
        "note": "This slide is a speed argument. The important shift is not replacement; it is throughput and iteration quality.",
    },
    {
        "layout": "split-list",
        "title": "Where AI Fits in Daily Work",
        "left_title": "Execution",
        "left": [
            "Coding and debugging",
            "Documentation and summaries",
            "Email, updates, and stakeholder communication",
        ],
        "right_title": "Decision Support",
        "right": [
            "Tradeoff analysis",
            "Planning and decomposition",
            "Requirement refinement",
        ],
        "note": "I keep this broad because the audience is mixed. The common pattern is converting context into structured output faster.",
    },
    {
        "layout": "cards",
        "title": "Tools I Use in Real Work",
        "cards": [
            ("ChatGPT", "Starting point", "Brainstorming, drafting, first-pass structure"),
            ("Claude", "Deep thinking", "Architecture, edge cases, long-context synthesis"),
            ("Codex", "Execution", "Repo-aware coding and task completion"),
            ("Copilot", "In-IDE speed", "Completions, tests, repetitive edits"),
        ],
        "note": "My daily pattern is stage-based. I switch tools based on the type of work, not loyalty to one tool.",
    },
    {
        "layout": "comparison",
        "title": "High-Level Positioning",
        "rows": [
            ("ChatGPT", "Best for starting", "Fast framing, broad help, quick drafts"),
            ("Claude", "Best for complexity", "Reasoning, architecture, long context"),
            ("Codex", "Best for execution", "Tool use, code changes, end-to-end tasks"),
            ("Copilot", "Best for velocity", "Inline coding acceleration inside the IDE"),
        ],
        "note": "The clean way to compare tools is by stage. Wrong tool selection creates extra prompt effort and cleanup work.",
    },
    {
        "layout": "flow",
        "title": "Use the Right Tool at the Right Stage",
        "steps": [
            ("1", "ChatGPT", "Frame the problem"),
            ("2", "Claude", "Refine the approach"),
            ("3", "Codex", "Execute the work"),
            ("4", "Copilot", "Accelerate local coding"),
        ],
        "note": "This is the workflow slide I actually use. The value comes from chaining tools intentionally rather than expecting one tool to do everything.",
    },
    {
        "layout": "numbered",
        "title": "How to Use AI",
        "subtitle": "Universal workflow for any role",
        "items": [
            ("01", "Define the goal clearly"),
            ("02", "Provide real context"),
            ("03", "Ask for a structured output"),
            ("04", "Review critically"),
            ("05", "Iterate in smaller loops"),
        ],
        "note": "This is the most reusable workflow in the deck. It improves both speed and output quality because it forces clarity first.",
    },
    {
        "layout": "prompt",
        "title": "Prompting Framework I Use Daily",
        "prompt_lines": [
            ("Act as", "the role you want the model to simulate"),
            ("Context", "business situation, stack, constraints, current state"),
            ("Task", "the exact job and expected outcome"),
            ("Output format", "table, bullets, code, slide, test plan, ADR"),
            ("Constraints", "style, security, depth, what not to do"),
        ],
        "note": "This is my default prompt skeleton. It scales across engineering, product, architecture, and documentation work.",
    },
    {
        "layout": "before-after",
        "title": "Prompt Quality Changes the Output",
        "before": "Write an API for leave requests.",
        "after": "Act as a Spring Boot engineer. Build a leave request API with validation, RBAC, paging, error handling, and tests. Return controller, DTOs, service flow, and test cases.",
        "footer": "Better prompts reduce cleanup.",
        "note": "One of the most common early mistakes is under-specifying the task and then blaming the tool for being generic.",
    },
    {
        "layout": "role",
        "title": "Developers",
        "highlight": "Where AI helps most",
        "bullets": [
            "Writing APIs, services, validation, and tests",
            "Explaining legacy code and unfamiliar modules",
            "Debugging stack traces and isolating likely faults",
            "Refactoring repetitive code safely with review",
        ],
        "note": "Developers usually see the biggest gains from generation plus explanation, not from blind code generation alone.",
    },
    {
        "layout": "code",
        "title": "Developer Workflow Example",
        "bullets": [
            "Prompt with stack, contract, validation rules, and response model",
            "Generate the API skeleton, then refine for auth and edge cases",
            "Ask for tests and then ask what is still missing",
        ],
        "code": [
            "@PostMapping(\"/leave-requests\")",
            "public LeaveResponse create(@Valid @RequestBody CreateLeaveRequest request) {",
            "    return leaveService.create(request);",
            "}",
        ],
        "note": "This workflow regularly saves meaningful implementation time, especially on repeatable service and controller work.",
    },
    {
        "layout": "role",
        "title": "Architects and Tech Leads",
        "highlight": "Where AI helps most",
        "bullets": [
            "Exploring options before committing to one design path",
            "Comparing tradeoffs across scale, cost, complexity, and risk",
            "Drafting HLDs, ADRs, sequences, and implementation slices",
            "Stress-testing assumptions by asking AI to argue against the design",
        ],
        "note": "Architect use is strongest when it is used to make assumptions explicit and alternatives visible.",
    },
    {
        "layout": "question-grid",
        "title": "Architect Example",
        "subtitle": "Prompt AI for a scalable system design",
        "questions": [
            "What are the core components?",
            "How does data flow through the system?",
            "Where are the bottlenecks at 10x scale?",
            "What alternatives should we evaluate?",
        ],
        "note": "I explicitly ask for bottlenecks and alternatives. That creates a better design discussion than asking for a single perfect answer.",
    },
    {
        "layout": "role",
        "title": "Product and Business Teams",
        "highlight": "Where AI helps most",
        "bullets": [
            "Turning rough ideas into PRDs and user stories",
            "Refining requirements before they hit engineering",
            "Summarizing meetings, calls, and stakeholder feedback",
            "Drafting updates, reports, and executive communication",
        ],
        "note": "For PM and business teams, AI is strongest when the input is messy and the output needs to be structured fast.",
    },
    {
        "layout": "flow-panel",
        "title": "PM and Business Workflow",
        "steps": [
            "Input a rough idea",
            "Generate a PRD",
            "Convert to stories and acceptance criteria",
            "Refine with the team",
            "Reuse AI for summaries and updates",
        ],
        "note": "This is one of the cleanest adoption paths because the artifacts are immediately useful to the wider team.",
    },
    {
        "layout": "demo",
        "title": "Live Demo #1",
        "subtitle": "Developer flow",
        "demo_title": "Build a REST API using AI",
        "bullets": [
            "Show prompt",
            "Inspect output",
            "Refine for validation and tests",
            "Explain what still needs engineering judgment",
        ],
        "note": "Keep the demo realistic. The value comes from showing prompt, output, critique, and refinement in sequence.",
    },
    {
        "layout": "demo",
        "title": "Live Demo #2",
        "subtitle": "Business and PM flow",
        "demo_title": "Turn raw notes into a structured summary",
        "bullets": [
            "Start with messy notes",
            "Extract decisions, risks, and actions",
            "Refine tone for stakeholders",
            "Show the before-to-after transformation",
        ],
        "note": "This usually lands with the full room because everyone understands the pain of converting notes into something useful.",
    },
    {
        "layout": "closing",
        "title": "Best Practices and Takeaways",
        "bullets": [
            "Treat AI like a junior engineer or analyst, not an authority",
            "Never skip review for code, decisions, data, or external claims",
            "Break large work into steps and reuse prompts that work well",
            "Start small, build habits fast, and scale by role",
        ],
        "quote": "AI is a productivity multiplier when paired with judgment.",
        "note": "I close by keeping the message grounded: practical, useful, and review-driven. That is how teams get durable value.",
    },
]


def add_full_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg"]

    top_glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-1.3), Inches(-2.2), Inches(6.0), Inches(5.0)
    )
    top_glow.fill.solid()
    top_glow.fill.fore_color.rgb = COLORS["blue"]
    top_glow.fill.transparency = 0.80
    top_glow.line.fill.background()

    right_glow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.8), Inches(-1.0), Inches(5.8), Inches(4.5)
    )
    right_glow.fill.solid()
    right_glow.fill.fore_color.rgb = COLORS["teal"]
    right_glow.fill.transparency = 0.82
    right_glow.line.fill.background()

    lower_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(6.45), SLIDE_W, Inches(1.05)
    )
    lower_panel.fill.solid()
    lower_panel.fill.fore_color.rgb = COLORS["bg_soft"]
    lower_panel.line.fill.background()

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.45), Inches(1.15), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["cyan"]
    line.line.fill.background()


def add_logo(slide, small=False):
    if LOGO_PATH.exists():
        width = Inches(1.55 if small else 1.95)
        slide.shapes.add_picture(str(LOGO_PATH), Inches(11.0), Inches(0.45), width=width)


def add_title(slide, title, subtitle=None, kicker=None):
    if kicker:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(0.72), Inches(3.5), Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = kicker
        r.font.name = "Aptos"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = COLORS["cyan"]

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.0), Inches(9.6), Inches(1.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = "Aptos Display"
    r.font.size = Pt(28 if len(title) > 26 else 32)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(2.02), Inches(8.5), Inches(0.65))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.name = "Aptos"
        sr.font.size = Pt(15)
        sr.font.color.rgb = COLORS["muted"]


def add_footer_index(slide, idx):
    box = slide.shapes.add_textbox(Inches(12.1), Inches(6.73), Inches(0.8), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{idx:02d}"
    r.font.name = "Aptos"
    r.font.size = Pt(11)
    r.font.color.rgb = COLORS["muted"]


def add_bullets(slide, bullets, x=0.85, y=2.35, w=7.2, h=3.5, size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.bullet = True
        p.space_after = Pt(11)
        r = p.add_run()
        r.text = bullet
        r.font.name = "Aptos"
        r.font.size = Pt(size)
        r.font.color.rgb = COLORS["text"]


def add_panel(slide, x, y, w, h, color="panel", radius=True):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS[color]
    panel.line.color.rgb = COLORS["line"]
    return panel


def render_title(slide, data):
    add_title(slide, data["title"], data.get("subtitle"), data.get("kicker"))
    quote = slide.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(7.0), Inches(1.2))
    tf = quote.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "How we actually use AI daily"
    r.font.name = "Aptos Display"
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = COLORS["cyan"]

    add_bullets(slide, [
        "Practical workflows for developers, architects, tech leads, and business teams",
        "Focused on speed, quality, and repeatable usage patterns",
        "Built around tools, prompts, demos, and role-specific scenarios",
    ], x=0.85, y=4.1, w=7.0, h=1.9, size=18)

    hero = add_panel(slide, 8.55, 1.25, 3.9, 4.8, color="panel_2")
    hero.line.color.rgb = COLORS["cyan"]
    for n, label in [("30-40%", "time saved on repeatable work"), ("4", "core tools used daily"), ("3", "roles covered live")]:
        pass
    stats = [("30-40%", "Typical savings on repeatable tasks"), ("4", "Tools used intentionally"), ("Today", "Actionable right after the session")]
    base_y = 1.75
    for i, (big, small) in enumerate(stats):
        bx = slide.shapes.add_textbox(Inches(8.9), Inches(base_y + i * 1.18), Inches(2.8), Inches(0.45))
        tf = bx.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = big
        r.font.name = "Aptos Display"
        r.font.size = Pt(25)
        r.font.bold = True
        r.font.color.rgb = COLORS["white"]

        sx = slide.shapes.add_textbox(Inches(8.9), Inches(base_y + 0.43 + i * 1.18), Inches(2.9), Inches(0.4))
        stf = sx.text_frame
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = small
        sr.font.name = "Aptos"
        sr.font.size = Pt(12)
        sr.font.color.rgb = COLORS["muted"]


def render_statement(slide, data):
    add_title(slide, data["title"])
    big = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.6), Inches(1.6))
    tf = big.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = data["statement"]
    r.font.name = "Aptos Display"
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]
    add_bullets(slide, data["bullets"], x=0.85, y=4.0, w=7.8, h=1.8, size=18)
    panel = add_panel(slide, 9.0, 2.0, 3.1, 2.4, color="panel_2")
    panel.line.color.rgb = COLORS["blue"]
    tx = slide.shapes.add_textbox(Inches(9.35), Inches(2.45), Inches(2.4), Inches(1.2))
    ttf = tx.text_frame
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    tr = tp.add_run()
    tr.text = "Faster cycles.\nBetter first drafts."
    tr.font.name = "Aptos Display"
    tr.font.size = Pt(22)
    tr.font.bold = True
    tr.font.color.rgb = COLORS["cyan"]


def render_split_list(slide, data):
    add_title(slide, data["title"])
    left = add_panel(slide, 0.8, 2.1, 5.7, 3.7)
    right = add_panel(slide, 6.8, 2.1, 5.7, 3.7, color="panel_2")
    for x, title, items in [(1.1, data["left_title"], data["left"]), (7.1, data["right_title"], data["right"])]:
        bx = slide.shapes.add_textbox(Inches(x), Inches(2.45), Inches(4.8), Inches(0.4))
        tf = bx.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = "Aptos Display"
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = COLORS["cyan"]
        add_bullets(slide, items, x=x, y=3.0, w=4.7, h=2.2, size=18)


def render_cards(slide, data):
    add_title(slide, data["title"])
    positions = [(0.8, 2.25), (6.7, 2.25), (0.8, 4.3), (6.7, 4.3)]
    colors = ["panel", "panel_2", "panel_2", "panel"]
    for (name, tag, desc), (x, y), color in zip(data["cards"], positions, colors):
        panel = add_panel(slide, x, y, 5.65, 1.65, color=color)
        panel.line.color.rgb = COLORS["line"]
        tagbox = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.28), Inches(1.3), Inches(0.25))
        tf = tagbox.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = tag.upper()
        r.font.name = "Aptos"
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = COLORS["amber"]
        nbox = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.55), Inches(2.3), Inches(0.35))
        ntf = nbox.text_frame
        np = ntf.paragraphs[0]
        nr = np.add_run()
        nr.text = name
        nr.font.name = "Aptos Display"
        nr.font.size = Pt(22)
        nr.font.bold = True
        nr.font.color.rgb = COLORS["white"]
        dbox = slide.shapes.add_textbox(Inches(x + 2.35), Inches(y + 0.54), Inches(2.85), Inches(0.5))
        dtf = dbox.text_frame
        dp = dtf.paragraphs[0]
        dr = dp.add_run()
        dr.text = desc
        dr.font.name = "Aptos"
        dr.font.size = Pt(13)
        dr.font.color.rgb = COLORS["muted"]


def render_comparison(slide, data):
    add_title(slide, data["title"])
    headers = [("Tool", 0.95), ("Positioning", 4.0), ("Use it for", 7.1)]
    for label, x in headers:
        box = slide.shapes.add_textbox(Inches(x), Inches(2.0), Inches(2.2), Inches(0.3))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = "Aptos"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = COLORS["cyan"]
    for i, row in enumerate(data["rows"]):
        y = 2.35 + i * 0.9
        add_panel(slide, 0.8, y, 11.7, 0.68, color="panel" if i % 2 == 0 else "panel_2")
        vals = [(row[0], 1.0, 18, True), (row[1], 4.05, 15, False), (row[2], 7.15, 14, False)]
        for text, x, size, bold in vals:
            box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.15), Inches(3.0 if x < 7 else 4.7), Inches(0.3))
            tf = box.text_frame
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = text
            r.font.name = "Aptos Display" if bold else "Aptos"
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = COLORS["white"] if bold else COLORS["text"]


def render_flow(slide, data):
    add_title(slide, data["title"])
    x = 0.9
    widths = [2.75, 2.75, 2.75, 2.75]
    colors = ["blue", "teal", "cyan", "amber"]
    for i, (num, tool, text) in enumerate(data["steps"]):
        add_panel(slide, x, 2.55, widths[i], 2.15, color="panel" if i % 2 == 0 else "panel_2")
        nb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(2.82), Inches(0.4), Inches(0.3))
        nt = nb.text_frame
        p = nt.paragraphs[0]
        r = p.add_run()
        r.text = num
        r.font.name = "Aptos Display"
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = COLORS[colors[i]]
        tb = slide.shapes.add_textbox(Inches(x + 0.25), Inches(3.2), Inches(2.0), Inches(0.4))
        ttf = tb.text_frame
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = tool
        tr.font.name = "Aptos Display"
        tr.font.size = Pt(24)
        tr.font.bold = True
        tr.font.color.rgb = COLORS["white"]
        db = slide.shapes.add_textbox(Inches(x + 0.25), Inches(3.75), Inches(2.2), Inches(0.45))
        dtf = db.text_frame
        dp = dtf.paragraphs[0]
        dr = dp.add_run()
        dr.text = text
        dr.font.name = "Aptos"
        dr.font.size = Pt(13)
        dr.font.color.rgb = COLORS["muted"]
        if i < 3:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 2.82), Inches(3.15), Inches(0.42), Inches(0.6))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["line"]
            arrow.line.fill.background()
        x += 3.0


def render_numbered(slide, data):
    add_title(slide, data["title"], data.get("subtitle"))
    for i, (num, label) in enumerate(data["items"]):
        y = 2.3 + i * 0.75
        nb = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(0.9), Inches(0.35))
        ntf = nb.text_frame
        np = ntf.paragraphs[0]
        nr = np.add_run()
        nr.text = num
        nr.font.name = "Aptos Display"
        nr.font.size = Pt(22)
        nr.font.bold = True
        nr.font.color.rgb = COLORS["cyan"]
        lb = slide.shapes.add_textbox(Inches(2.0), Inches(y + 0.02), Inches(5.8), Inches(0.35))
        ltf = lb.text_frame
        lp = ltf.paragraphs[0]
        lr = lp.add_run()
        lr.text = label
        lr.font.name = "Aptos"
        lr.font.size = Pt(20)
        lr.font.color.rgb = COLORS["text"]
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.65), Inches(y + 0.1), Inches(2.7 - i * 0.35), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS["blue"] if i < 2 else COLORS["teal"] if i < 4 else COLORS["amber"]
        line.line.fill.background()


def render_prompt(slide, data):
    add_title(slide, data["title"])
    panel = add_panel(slide, 0.85, 2.15, 11.4, 3.65, color="panel_2")
    panel.line.color.rgb = COLORS["cyan"]
    for i, (head, desc) in enumerate(data["prompt_lines"]):
        y = 2.5 + i * 0.64
        hb = slide.shapes.add_textbox(Inches(1.15), Inches(y), Inches(1.7), Inches(0.3))
        htf = hb.text_frame
        hp = htf.paragraphs[0]
        hr = hp.add_run()
        hr.text = head
        hr.font.name = "Aptos Display"
        hr.font.size = Pt(18)
        hr.font.bold = True
        hr.font.color.rgb = COLORS["amber"]
        db = slide.shapes.add_textbox(Inches(3.0), Inches(y), Inches(7.8), Inches(0.32))
        dtf = db.text_frame
        dp = dtf.paragraphs[0]
        dr = dp.add_run()
        dr.text = desc
        dr.font.name = "Aptos"
        dr.font.size = Pt(17)
        dr.font.color.rgb = COLORS["text"]


def render_before_after(slide, data):
    add_title(slide, data["title"])
    left = add_panel(slide, 0.95, 2.35, 5.45, 2.8, color="panel")
    right = add_panel(slide, 6.9, 2.35, 5.45, 2.8, color="panel_2")
    left.line.color.rgb = COLORS["rose"]
    right.line.color.rgb = COLORS["green"]
    for x, label, text, color in [(1.2, "Before", data["before"], "rose"), (7.15, "After", data["after"], "green")]:
        hb = slide.shapes.add_textbox(Inches(x), Inches(2.7), Inches(1.0), Inches(0.3))
        htf = hb.text_frame
        hp = htf.paragraphs[0]
        hr = hp.add_run()
        hr.text = label
        hr.font.name = "Aptos"
        hr.font.size = Pt(11)
        hr.font.bold = True
        hr.font.color.rgb = COLORS[color]
        tb = slide.shapes.add_textbox(Inches(x), Inches(3.1), Inches(4.9), Inches(1.6))
        ttf = tb.text_frame
        tp = ttf.paragraphs[0]
        tr = tp.add_run()
        tr.text = text
        tr.font.name = "Aptos"
        tr.font.size = Pt(17 if label == "After" else 18)
        tr.font.color.rgb = COLORS["text"]
    fb = slide.shapes.add_textbox(Inches(0.95), Inches(5.45), Inches(4.0), Inches(0.35))
    ftf = fb.text_frame
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    fr.text = data["footer"]
    fr.font.name = "Aptos Display"
    fr.font.size = Pt(18)
    fr.font.bold = True
    fr.font.color.rgb = COLORS["cyan"]


def render_role(slide, data):
    add_title(slide, data["title"])
    hb = slide.shapes.add_textbox(Inches(0.85), Inches(2.0), Inches(4.5), Inches(0.5))
    htf = hb.text_frame
    hp = htf.paragraphs[0]
    hr = hp.add_run()
    hr.text = data["highlight"]
    hr.font.name = "Aptos Display"
    hr.font.size = Pt(26)
    hr.font.bold = True
    hr.font.color.rgb = COLORS["cyan"]
    add_bullets(slide, data["bullets"], x=0.9, y=2.85, w=6.9, h=2.8, size=19)
    stat = add_panel(slide, 8.5, 2.1, 3.8, 3.2, color="panel_2")
    stat.line.color.rgb = COLORS["line"]
    tb = slide.shapes.add_textbox(Inches(8.85), Inches(2.65), Inches(3.0), Inches(1.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Practical usage\nbeats generic prompting."
    r.font.name = "Aptos Display"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]


def render_code(slide, data):
    add_title(slide, data["title"])
    add_bullets(slide, data["bullets"], x=0.9, y=2.35, w=5.6, h=2.6, size=18)
    code_panel = add_panel(slide, 6.8, 2.2, 5.0, 3.4, color="panel_2")
    code_panel.line.color.rgb = COLORS["cyan"]
    mac = [COLORS["rose"], COLORS["amber"], COLORS["green"]]
    for i, color in enumerate(mac):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(7.1 + i * 0.24), Inches(2.45), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
    box = slide.shapes.add_textbox(Inches(7.1), Inches(2.85), Inches(4.3), Inches(2.2))
    tf = box.text_frame
    for i, line in enumerate(data["code"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        r = p.add_run()
        r.text = line
        r.font.name = "Courier New"
        r.font.size = Pt(13)
        r.font.color.rgb = COLORS["white"]


def render_question_grid(slide, data):
    add_title(slide, data["title"], data.get("subtitle"))
    positions = [(0.9, 2.5), (6.7, 2.5), (0.9, 4.25), (6.7, 4.25)]
    for q, (x, y) in zip(data["questions"], positions):
        add_panel(slide, x, y, 5.05, 1.25, color="panel" if x < 2 else "panel_2")
        box = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 0.28), Inches(4.45), Inches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = q
        r.font.name = "Aptos Display"
        r.font.size = Pt(19)
        r.font.bold = True
        r.font.color.rgb = COLORS["white"]


def render_flow_panel(slide, data):
    add_title(slide, data["title"])
    main = add_panel(slide, 0.95, 2.15, 11.25, 3.7, color="panel")
    main.line.color.rgb = COLORS["line"]
    for i, step in enumerate(data["steps"]):
        x = 1.35 + i * 2.08
        num = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(3.05), Inches(0.42), Inches(0.42))
        num.fill.solid()
        num.fill.fore_color.rgb = COLORS["cyan"] if i < 2 else COLORS["teal"] if i < 4 else COLORS["amber"]
        num.line.fill.background()
        box = slide.shapes.add_textbox(Inches(x + 0.58), Inches(3.0), Inches(1.3), Inches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = step
        r.font.name = "Aptos"
        r.font.size = Pt(16)
        r.font.color.rgb = COLORS["text"]
        if i < len(data["steps"]) - 1:
            ch = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x + 1.55), Inches(3.03), Inches(0.35), Inches(0.45))
            ch.fill.solid()
            ch.fill.fore_color.rgb = COLORS["line"]
            ch.line.fill.background()


def render_demo(slide, data):
    add_title(slide, data["title"], data.get("subtitle"), "LIVE")
    demo = add_panel(slide, 0.9, 2.0, 7.25, 3.95, color="panel_2")
    demo.line.color.rgb = COLORS["cyan"]
    tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.4), Inches(5.8), Inches(0.45))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = data["demo_title"]
    r.font.name = "Aptos Display"
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = COLORS["white"]
    window = add_panel(slide, 8.55, 2.05, 3.55, 3.9, color="panel")
    window.line.color.rgb = COLORS["line"]
    for i, color in enumerate([COLORS["rose"], COLORS["amber"], COLORS["green"]]):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.85 + i * 0.22), Inches(2.3), Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
    add_bullets(slide, data["bullets"], x=1.15, y=3.25, w=5.8, h=2.0, size=18)
    sb = slide.shapes.add_textbox(Inches(8.9), Inches(2.75), Inches(2.7), Inches(2.2))
    stf = sb.text_frame
    for i, line in enumerate(["Prompt", "Output", "Review", "Refine"]):
        p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
        p.space_after = Pt(12)
        r = p.add_run()
        r.text = line
        r.font.name = "Aptos Display"
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = COLORS["cyan"] if i % 2 == 0 else COLORS["text"]


def render_closing(slide, data):
    add_title(slide, data["title"])
    quote = slide.shapes.add_textbox(Inches(0.85), Inches(2.05), Inches(7.6), Inches(0.9))
    tf = quote.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = data["quote"]
    r.font.name = "Aptos Display"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = COLORS["cyan"]
    add_bullets(slide, data["bullets"], x=0.9, y=3.15, w=7.3, h=2.6, size=18)
    panel = add_panel(slide, 8.7, 2.2, 3.45, 3.5, color="panel_2")
    panel.line.color.rgb = COLORS["amber"]
    bx = slide.shapes.add_textbox(Inches(9.05), Inches(3.0), Inches(2.7), Inches(1.2))
    btf = bx.text_frame
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = "Start small.\nScale by role.\nKeep standards high."
    br.font.name = "Aptos Display"
    br.font.size = Pt(22)
    br.font.bold = True
    br.font.color.rgb = COLORS["white"]


LAYOUTS = {
    "title": render_title,
    "statement": render_statement,
    "split-list": render_split_list,
    "cards": render_cards,
    "comparison": render_comparison,
    "flow": render_flow,
    "numbered": render_numbered,
    "prompt": render_prompt,
    "before-after": render_before_after,
    "role": render_role,
    "code": render_code,
    "question-grid": render_question_grid,
    "flow-panel": render_flow_panel,
    "demo": render_demo,
    "closing": render_closing,
}


def add_notes(slide, note_text):
    slide.notes_slide.notes_text_frame.text = note_text


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for idx, data in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        add_full_bg(slide)
        add_logo(slide, small=(data["layout"] != "title"))
        LAYOUTS[data["layout"]](slide, data)
        add_footer_index(slide, idx)
        add_notes(slide, data["note"])

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))


if __name__ == "__main__":
    build_presentation()
